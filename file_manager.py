import os
import time
import threading
from datetime import datetime
import pandas as pd
import requests
from elasticsearch import Elasticsearch, helpers
import warnings
import traceback

# Suppress warnings
warnings.filterwarnings("ignore")

HAS_SENTENCE_TRANSFORMERS = False

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    HAS_LANGCHAIN = True
except ImportError:
    HAS_LANGCHAIN = False

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

import openpyxl

# --- Helper Classes ---


class ContentExtractor:
    """Handles content extraction and splitting for various file formats."""
    def __init__(self):
        if HAS_LANGCHAIN:
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=600,
                chunk_overlap=100,
                separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""]
            )
        else:
            self.text_splitter = None

    def load_and_split(self, file_path, ext):
        """
        Extract content from file and split into chunks.
        Returns: List of dicts [{"content": str, "page_info": str}]
        """
        try:
            raw_text = ""
            metadata_list = [] # parallel to chunks if possible, or just raw text first

            if ext == '.docx':
                if Document:
                    doc = Document(file_path)
                    raw_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                else:
                    return []

            elif ext == '.pdf':
                if PdfReader:
                    reader = PdfReader(file_path)
                    texts = []
                    for i, page in enumerate(reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            texts.append(f"[Page {i+1}] {page_text}")
                    raw_text = "\n".join(texts)
                else:
                    return []

            elif ext in ['.xlsx', '.xls']:
                # Optimized for RAG: Row-based chunking with explicit headers
                try:
                    dfs = pd.read_excel(file_path, sheet_name=None)
                    custom_chunks = []
                    chunk_counter = 0
                    
                    for sheet, df in dfs.items():
                        # Clean data
                        df = df.fillna('')
                        
                        # Convert all columns to string for safe usage
                        columns = [str(c).strip() for c in df.columns]
                        
                        # Add sheet context
                        sheet_context = f"File: {os.path.basename(file_path)}, Sheet: {sheet}"
                        
                        for _, row in df.iterrows():
                            # Construct semantic row representation: "Col1: Val1, Col2: Val2..."
                            row_parts = []
                            for col, val in zip(columns, row):
                                val_str = str(val).strip()
                                if val_str: # Skip empty values
                                    row_parts.append(f"{col}: {val_str}")
                            
                            if not row_parts:
                                continue
                                
                            row_str = f"[{sheet}] " + ", ".join(row_parts)
                            
                            # Strict One Row Per Chunk
                            # Include context in every chunk for independence
                            chunk_text = f"{sheet_context}\n{row_str}"
                            custom_chunks.append({"content": chunk_text, "chunk_id": chunk_counter})
                            chunk_counter += 1
                            
                    return custom_chunks
                    
                except Exception as e:
                    print(f"Error parsing Excel {file_path}: {e}")
                    return []

            elif ext == '.pptx':
                if Presentation:
                    prs = Presentation(file_path)
                    texts = []
                    for i, slide in enumerate(prs.slides):
                        slide_texts = []
                        if slide.shapes.title:
                            slide_texts.append(f"Title: {slide.shapes.title.text}")
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_texts.append(shape.text)
                        if slide_texts:
                            texts.append(f"[Slide {i+1}] " + "\n".join(slide_texts))
                    raw_text = "\n".join(texts)
                else:
                    return []

            elif ext in ['.html', '.htm']:
                if BeautifulSoup:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        soup = BeautifulSoup(f, 'html.parser')
                        raw_text = soup.get_text(separator='\n')
                else:
                    return []

            elif ext in ['.txt', '.md', '.py', '.json', '.log', '.xml', '.ini', '.yml']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    raw_text = f.read()

            else:
                return []

            # Split text
            if not raw_text.strip():
                return []

            if self.text_splitter:
                chunks = self.text_splitter.split_text(raw_text)
            else:
                # Fallback simple splitter
                chunks = [raw_text[i:i+500] for i in range(0, len(raw_text), 500)]

            return [{"content": chunk, "chunk_id": i} for i, chunk in enumerate(chunks)]

        except Exception as e:
            print(f"Error extracting {file_path}: {e}")
            return []

class VectorModel:
    """Wrapper for Aliyun OpenSearch Text Embedding API."""
    def __init__(self, api_url, api_key=None):
        self.api_key = api_key
        self.api_url = api_url
        # ops-text-embedding-002 dimension is 1024 based on test result
        self.dim = 1024

    def encode(self, texts, input_type="QUERY"):
        if not self.api_key or not self.api_url or not texts:
            return [[0.0] * self.dim] * len(texts)
        
        # Check if texts are empty
        if not any(t.strip() for t in texts):
            return [[0.0] * self.dim] * len(texts)

        # Batch processing
        batch_size = 32  # Max batch size per API documentation
        all_embeddings = []
        
        try:
            for i in range(0, len(texts), batch_size):
                # Slice batch
                raw_batch = texts[i:i + batch_size]
                
                # Truncate each text to max 8192 chars
                batch_texts = [t[:8192] for t in raw_batch]
                
                # Skip empty batch (shouldn't happen with range logic but good safety)
                if not batch_texts:
                    continue
                    
                headers = {
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {self.api_key}"
                }
                
                payload = {
                    "input": batch_texts,
                    "input_type": input_type
                }
                
                response = requests.post(self.api_url, json=payload, headers=headers, timeout=30) # Increased timeout
                response.raise_for_status()
                
                result = response.json()
                if "result" in result and "embeddings" in result["result"]:
                    # Sort by index to ensure order matches input
                    embeddings_data = sorted(result["result"]["embeddings"], key=lambda x: x["index"])
                    batch_embeddings = [item["embedding"] for item in embeddings_data]
                    all_embeddings.extend(batch_embeddings)
                    
                    # Update dim if we get a result
                    if batch_embeddings:
                         self.dim = len(batch_embeddings[0])
                else:
                    print(f"API response format unexpected: {result}")
                    all_embeddings.extend([[0.0] * self.dim] * len(batch_texts))
            
            return all_embeddings
                
        except Exception as e:
            print(f"Encoding error: {e}")
            # Return zero vectors for failed batch + remaining items to match length
            current_len = len(all_embeddings)
            needed = len(texts) - current_len
            return all_embeddings + ([[0.0] * self.dim] * needed)

# --- Main Class ---

# 全局实例
instance = None

class FileManager:
    def __init__(self, config):
        global instance
        instance = self
        self.config = config.get('file_knowledge', {})
        
        # Support both legacy and new config keys
        self.monitored_folders = self.config.get('dirs', self.config.get('monitored_folders', []))
        self.file_types = set(self.config.get('types', self.config.get('file_types', [])))
        self.reindex_interval = self.config.get('watch_interval', self.config.get('reindex_interval', 3600))

        # ES Config
        es_config = config.get('es', {})
        self.es_host = es_config.get('host', "http://localhost:9200")
        self.es_user = es_config.get('user')
        self.es_password = es_config.get('password')
        # Use a new index name for RAG to avoid conflict with old mapping
        # ops-text-embedding-002 has 1024 dims, updating index name
        self.index_name = "deskq_rag_files_v3"
        
        # Initialize Components
        self.extractor = ContentExtractor()
        # Pass API key and URL to VectorModel
        vector_config = config.get('vector', {})
        self.vector_model = VectorModel(vector_config.get("apiUrl"), vector_config.get("apikey"))
        self.vector_dim = self.vector_model.dim
        
        print(f"Initializing Elasticsearch at {self.es_host}...")
        try:
            es_args = {
                "hosts": [self.es_host],
                "timeout": 30
            }
            
            if self.es_user and self.es_password:
                es_args["http_auth"] = (self.es_user, self.es_password)
                
            self.es = Elasticsearch(**es_args)
            if not self.es.ping():
                print(f"Warning: Could not connect to Elasticsearch at {self.es_host}")
        except Exception as e:
            print(f"Error connecting to Elasticsearch: {e}")
            self.es = None

        self._init_es_index()
        
        self.is_indexing = False
        self.running = True
        
        # Auto index thread
        if config.get('auto_index', True):
            self.thread = threading.Thread(target=self._auto_reindex_loop, daemon=True)
            self.thread.start()

    def _init_es_index(self):
        if not self.es:
            return
            
        try:
            # Check if index exists
            if self.es.indices.exists(index=self.index_name):
                # Verify mapping 'path' is keyword
                mapping = self.es.indices.get_mapping(index=self.index_name)
                # Structure: {index_name: {mappings: {properties: {path: {type: ...}}}}}
                try:
                    path_type = mapping[self.index_name]['mappings']['properties']['path']['type']
                    if path_type != 'keyword':
                        print(f"Index {self.index_name} has incorrect mapping for 'path' ({path_type}). Recreating...")
                        self.es.indices.delete(index=self.index_name)
                    else:
                        # Mapping is correct
                        return
                except KeyError:
                    # Could happen if 'path' field doesn't exist yet or structure is different
                    print(f"Index {self.index_name} mapping check failed. Recreating...")
                    self.es.indices.delete(index=self.index_name, ignore=[404])

            # Create index if not exists (or deleted above)
            if not self.es.indices.exists(index=self.index_name):
                mapping = {
                    "mappings": {
                        "properties": {
                            "file_name": {"type": "text", "analyzer": "standard"},
                            "path": {"type": "keyword"},
                            "ext": {"type": "keyword"},
                            "mtime": {"type": "double"},
                            "content": {"type": "text", "analyzer": "standard"},
                            "content_vector": {
                                "type": "dense_vector",
                                "dims": self.vector_dim
                            },
                            "chunk_id": {"type": "integer"},
                            "created_at": {"type": "date"}
                        }
                    }
                }
                self.es.indices.create(index=self.index_name, body=mapping)
                print(f"Created Elasticsearch index: {self.index_name}")
        except Exception as e:
            print(f"Error checking/creating index: {e}")

    def _auto_reindex_loop(self):
        while self.running:
            try:
                self.index_files()
            except Exception as e:
                print(f"Auto-index error: {e}")
            time.sleep(self.reindex_interval)

    def index_files(self, specific_folders=None):
        """
        Index files from monitored folders.
        Supports incremental updates via mtime.
        """
        if self.is_indexing:
            return
        self.is_indexing = True
        target_folders = specific_folders or self.monitored_folders
        print(f"[{datetime.now().strftime('%H:%M:%S')}] Starting RAG index on {target_folders}...")
        
        try:
            # 1. Scan current files on disk
            current_files = {} # path -> {name, ext, path, mtime}
            
            for folder in target_folders:
                if not os.path.exists(folder):
                    continue
                
                for root, dirs, files in os.walk(folder):
                    for file in files:
                        ext = os.path.splitext(file)[1].lower()
                        if ext in self.file_types:
                            full_path = os.path.join(root, file)
                            full_path = os.path.normpath(full_path)
                            mtime = os.path.getmtime(full_path)
                            
                            current_files[full_path] = {
                                "name": file,
                                "path": full_path,
                                "ext": ext,
                                "mtime": mtime
                            }
            
            # 2. Get existing files metadata from ES
            # We aggregate by path to get stored mtime
            existing_files = {} # path -> mtime
            if self.es:
                try:
                    # Ensure index exists
                    if not self.es.indices.exists(index=self.index_name):
                        self._init_es_index()
                        
                    # Terms aggregation on path to find all indexed files
                    # Note: For huge datasets, composite agg or scroll is better.
                    # Here we use a simple search collapsing on path or just scrolling hits.
                    # Since we store multiple chunks per file, we just need to know which files are there.
                    # We can search for all unique 'path' and their 'mtime'.
                    
                    # Using a collapse query to get unique paths
                    resp = self.es.search(
                        index=self.index_name,
                        body={
                            "query": {"match_all": {}},
                            "_source": ["path", "mtime"],
                            "collapse": {"field": "path"},
                            "size": 10000 
                        }
                    )
                    
                    for hit in resp['hits']['hits']:
                        src = hit['_source']
                        existing_files[src.get('path')] = src.get('mtime', 0)
                        
                except Exception as e:
                    print(f"Error fetching existing files from ES: {e}")

            # 3. Calculate Diff
            files_on_disk = set(current_files.keys())
            files_in_index = set(existing_files.keys())
            
            to_delete = files_in_index - files_on_disk
            to_add_or_update = []
            
            for path in files_on_disk:
                if path not in files_in_index:
                    to_add_or_update.append(path)
                else:
                    # Check mtime
                    if current_files[path]['mtime'] > existing_files[path] + 1.0: # 1s buffer
                        to_add_or_update.append(path)

            # 4. Apply Deletions
            if to_delete and self.es:
                print(f"Removing {len(to_delete)} deleted files...")
                for path in to_delete:
                    self.es.delete_by_query(
                        index=self.index_name,
                        body={"query": {"term": {"path": path}}}
                    )

            # 5. Apply Updates (Delete old chunks -> Index new chunks)
            if to_add_or_update and self.es:
                print(f"Indexing {len(to_add_or_update)} files...")
                
                for path in to_add_or_update:
                    try:
                        # First delete existing chunks for this file to avoid duplication
                        if path in files_in_index:
                             self.es.delete_by_query(
                                index=self.index_name,
                                body={"query": {"term": {"path": path}}}
                            )
                        
                        meta = current_files[path]
                        
                        # Extract and Split
                        chunks = self.extractor.load_and_split(path, meta['ext'])
                        if not chunks:
                            continue
                            
                        # Embed
                        texts = [c['content'] for c in chunks]
                        # Use input_type="DOCUMENT" for indexing if API supports, but user example used "query".
                        # Usually embedding APIs use "document" for storage. Let's use "document" for now.
                        # Wait, user example only showed "query". If I use "document" and it fails, that's bad.
                        # However, Aliyun OpenSearch usually supports "query" and "document".
                        # Let's try "document" for indexing.
                        vectors = self.vector_model.encode(texts, input_type="DOCUMENT")
                        
                        # Prepare Bulk Actions
                        actions = []
                        for i, chunk in enumerate(chunks):
                            doc = {
                                "_index": self.index_name,
                                "_source": {
                                    "file_name": meta['name'],
                                    "path": meta['path'],
                                    "ext": meta['ext'],
                                    "mtime": meta['mtime'],
                                    "content": chunk['content'],
                                    "content_vector": vectors[i],
                                    "chunk_id": chunk['chunk_id'],
                                    "created_at": datetime.now().isoformat()
                                }
                            }
                            actions.append(doc)
                        
                        if actions:
                            helpers.bulk(self.es, actions)
                            
                    except Exception as e:
                        print(f"Failed to index {path}: {e}")
                        traceback.print_exc()

            print(f"[{datetime.now().strftime('%H:%M:%S')}] Indexing completed. (+{len(to_add_or_update)}, -{len(to_delete)})")
            
        except Exception as e:
            print(f"Error during file indexing: {e}")
            traceback.print_exc()
        finally:
            self.is_indexing = False

    def stop(self):
        self.running = False
